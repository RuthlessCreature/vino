from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
MD_PATH = ROOT / "docs" / "vino-boss-deck.md"
PPTX_PATH = ROOT / "docs" / "vino-boss-deck.pptx"

BG = RGBColor(0x05, 0x06, 0x08)
PANEL = RGBColor(0x0C, 0x10, 0x14)
ACCENT = RGBColor(0x62, 0xF0, 0xFF)
TEXT = RGBColor(0xF3, 0xF6, 0xF8)
MUTED = RGBColor(0xA8, 0xB7, 0xC2)
SUCCESS = RGBColor(0x55, 0xE3, 0x9E)
LINE = RGBColor(0x24, 0x30, 0x3A)


def parse_markdown(path: Path):
    text = path.read_text(encoding="utf-8")
    lines = text.splitlines()

    title = next((line[2:].strip() for line in lines if line.startswith("# ")), "vino 商业逻辑汇报")

    subtitle_lines = []
    for line in lines:
        if line.startswith("> "):
            subtitle_lines.append(line[2:].strip())
        elif subtitle_lines and not line.strip():
            continue
        elif subtitle_lines:
            break
    subtitle = " | ".join(subtitle_lines)

    sections = []
    current = None
    for line in lines:
        if line.startswith("## "):
            if current:
                sections.append(current)
            current = {"title": line[3:].strip(), "lines": []}
        elif current is not None:
            current["lines"].append(line.rstrip())
    if current:
        sections.append(current)

    return title, subtitle, sections


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_header_bar(slide, title):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.45),
        Inches(0.35),
        Inches(12.25),
        Inches(0.7),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = LINE

    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.margin_left = Pt(16)
    text_frame.margin_top = Pt(10)

    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = TEXT
    paragraph.alignment = PP_ALIGN.LEFT


def add_footer(slide, index, total):
    box = slide.shapes.add_textbox(Inches(11.7), Inches(7.0), Inches(1.0), Inches(0.25))
    text_frame = box.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = f"{index}/{total}"
    run.font.size = Pt(11)
    run.font.color.rgb = MUTED
    paragraph.alignment = PP_ALIGN.RIGHT


def add_content_box(slide):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.65),
        Inches(1.35),
        Inches(12.0),
        Inches(5.25),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.fill.transparency = 0.05
    shape.line.color.rgb = LINE
    return shape


def add_line(text_frame, text, *, first, level=0, bullet=False, accent=False, bold=False):
    paragraph = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
    paragraph.text = text
    paragraph.level = level
    paragraph.font.size = Pt(24 if bullet and level == 0 else 20 if bullet else 22)
    paragraph.font.color.rgb = ACCENT if accent else (TEXT if level == 0 else MUTED)
    paragraph.font.bold = bold
    paragraph.space_after = Pt(8 if bullet else 10)
    if bullet:
        paragraph.bullet = True


def fill_text_frame(text_frame, content_lines):
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.margin_left = Pt(18)
    text_frame.margin_right = Pt(16)
    text_frame.margin_top = Pt(16)
    text_frame.margin_bottom = Pt(12)

    first = True
    for raw in content_lines:
        line = raw.rstrip()
        if not line.strip():
            continue
        stripped = line.lstrip(" ")
        if stripped.startswith("- "):
            indent = len(line) - len(stripped)
            level = min(indent // 2, 2)
            add_line(text_frame, stripped[2:].strip(), first=first, level=level, bullet=True)
        else:
            accent = line.strip().endswith("：")
            add_line(text_frame, line.strip(), first=first, accent=accent, bold=accent)
        first = False


def add_cover_slide(prs: Presentation, title: str, subtitle: str, section: dict, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    left.fill.solid()
    left.fill.fore_color.rgb = ACCENT
    left.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(10.8), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.clear()
    paragraph = title_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = title
    run.font.size = Pt(31)
    run.font.bold = True
    run.font.color.rgb = TEXT

    subtitle_box = slide.shapes.add_textbox(Inches(1.02), Inches(1.95), Inches(10.8), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.clear()
    subtitle_paragraph = subtitle_frame.paragraphs[0]
    subtitle_run = subtitle_paragraph.add_run()
    subtitle_run.text = subtitle
    subtitle_run.font.size = Pt(15)
    subtitle_run.font.color.rgb = MUTED

    panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.95),
        Inches(3.0),
        Inches(11.3),
        Inches(2.7),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = PANEL
    panel.line.color.rgb = LINE
    fill_text_frame(panel.text_frame, section["lines"])

    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(10.2),
        Inches(6.28),
        Inches(1.95),
        Inches(0.42),
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = SUCCESS
    badge.line.fill.background()
    badge_frame = badge.text_frame
    badge_frame.clear()
    badge_paragraph = badge_frame.paragraphs[0]
    badge_paragraph.alignment = PP_ALIGN.CENTER
    badge_run = badge_paragraph.add_run()
    badge_run.text = "Boss Deck"
    badge_run.font.size = Pt(14)
    badge_run.font.bold = True
    badge_run.font.color.rgb = BG

    add_footer(slide, 1, total)


def add_section_slide(prs: Presentation, title: str, content_lines: list[str], index: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header_bar(slide, title)
    box = add_content_box(slide)
    fill_text_frame(box.text_frame, content_lines)
    add_footer(slide, index, total)


def build_presentation():
    title, subtitle, sections = parse_markdown(MD_PATH)
    if not sections:
        raise RuntimeError("No sections found in markdown source.")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    total = len(sections)
    add_cover_slide(prs, title, subtitle, sections[0], total)

    for index, section in enumerate(sections[1:], start=2):
        add_section_slide(prs, section["title"], section["lines"], index, total)

    prs.save(PPTX_PATH)
    return PPTX_PATH


if __name__ == "__main__":
    path = build_presentation()
    print(path)
