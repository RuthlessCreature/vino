# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path(__file__).with_name("iphone-model-download-ux-flow.pptx")


class C:
    bg = RGBColor(248, 250, 252)
    panel = RGBColor(255, 255, 255)
    ink = RGBColor(15, 23, 42)
    muted = RGBColor(71, 85, 105)
    line = RGBColor(203, 213, 225)
    green = RGBColor(34, 197, 94)
    blue = RGBColor(37, 99, 235)
    cyan = RGBColor(8, 145, 178)
    amber = RGBColor(245, 158, 11)
    red = RGBColor(220, 38, 38)
    slate = RGBColor(51, 65, 85)
    dark = RGBColor(2, 6, 23)
    white = RGBColor(255, 255, 255)


FONT = "Microsoft YaHei"


def prs_new():
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    return prs


def set_bg(slide, color=C.bg):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, x, y, w, h, text="", size=18, color=C.ink, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.62, 0.35, 11.7, 0.5, title, size=25, bold=True, color=C.ink)
    if subtitle:
        add_textbox(slide, 0.64, 0.88, 11.9, 0.42, subtitle, size=11, color=C.muted)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.64), Inches(1.22), Inches(12.05), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = C.line
    line.line.fill.background()


def add_label(slide, x, y, w, h, text, color=C.blue, text_color=C.white):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = text_color
    return shape


def add_card(slide, x, y, w, h, title, body=None, accent=C.blue, fill=C.panel):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = C.line
    shape.line.width = Pt(1)

    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()

    add_textbox(slide, x + 0.18, y + 0.16, w - 0.34, 0.33, title, size=14, bold=True, color=C.ink)
    if body:
        add_textbox(slide, x + 0.18, y + 0.56, w - 0.34, h - 0.66, body, size=10, color=C.muted)
    return shape


def add_bullets(slide, x, y, w, h, items, size=13, color=C.ink, leading=1.15):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(5 * leading)
    return box


def add_numbered(slide, x, y, w, h, items, size=12, color=C.ink):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items, start=1):
        p = tf.paragraphs[0] if i == 1 else tf.add_paragraph()
        p.text = f"{i}. {item}"
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(4)
    return box


def add_flow(slide, x, y, labels, colors=None, box_w=1.55, box_h=0.64, gap=0.22):
    colors = colors or [C.blue] * len(labels)
    cursor = x
    for idx, label in enumerate(labels):
        add_label(slide, cursor, y, box_w, box_h, label, color=colors[idx])
        cursor += box_w
        if idx != len(labels) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(cursor + 0.02), Inches(y + 0.16), Inches(gap - 0.04), Inches(0.32))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = C.line
            arrow.line.fill.background()
            cursor += gap


def add_table_like(slide, x, y, rows, widths, row_h=0.46, header=False):
    top = y
    for r_idx, row in enumerate(rows):
        left = x
        for c_idx, text in enumerate(row):
            color = C.slate if header and r_idx == 0 else C.panel
            text_color = C.white if header and r_idx == 0 else C.ink
            rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(widths[c_idx]), Inches(row_h))
            rect.fill.solid()
            rect.fill.fore_color.rgb = color
            rect.line.color.rgb = C.line
            tf = rect.text_frame
            tf.clear()
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            rr = p.add_run()
            rr.text = text
            rr.font.name = FONT
            rr.font.size = Pt(9.5)
            rr.font.bold = bool(header and r_idx == 0)
            rr.font.color.rgb = text_color
            left += widths[c_idx]
        top += row_h


def build():
    prs = prs_new()

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C.dark)
    add_textbox(slide, 0.72, 0.75, 11.7, 0.95, "vino_iPhone 模型下载体验升级", size=33, bold=True, color=C.white)
    add_textbox(slide, 0.76, 1.75, 11.0, 0.55, "产品形态与操作流程：从手动 URL + 账号密码，升级为平台驱动的一次性设备绑定", size=17, color=RGBColor(203, 213, 225))
    add_flow(
        slide,
        0.82,
        3.05,
        ["打开 App", "绑定平台", "同步模型", "下载安装", "模型就绪"],
        colors=[C.blue, C.cyan, C.green, C.amber, C.green],
        box_w=1.75,
        box_h=0.7,
        gap=0.38,
    )
    add_card(slide, 0.82, 4.55, 3.65, 1.25, "核心目标", "普通用户不需要理解服务器地址、token、下载票据，只需要确认设备属于哪个平台账号。", accent=C.green, fill=RGBColor(15, 23, 42))
    add_card(slide, 4.82, 4.55, 3.65, 1.25, "产品边界", "保留账号密码和 URL 输入作为高级诊断入口，不放在日常主路径。", accent=C.amber, fill=RGBColor(15, 23, 42))
    add_card(slide, 8.82, 4.55, 3.65, 1.25, "技术原则", "沿用现有 /api/cloud/v1 模型清单、download-ticket、加密下载、设备绑定链路。", accent=C.blue, fill=RGBColor(15, 23, 42))

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "当前状态：下载能力有了，但首次体验偏工程化", "代码现状来自 vino_iPhone 当前工程；问题集中在入口设计，不是下载协议本身。")
    add_flow(slide, 0.72, 1.62, ["输入 URL", "输入账号密码", "登录", "同步清单", "选择模型", "下载安装"], colors=[C.red, C.red, C.blue, C.cyan, C.amber, C.green], box_w=1.45, box_h=0.58, gap=0.28)
    add_card(slide, 0.72, 2.65, 3.7, 3.55, "用户痛点", "1. URL 对终端用户没有产品语义。\n2. 真机上 127.0.0.1 默认地址天然不可用。\n3. 相机 Overlay 暴露账号、密码、网址，干扰主工作流。\n4. token 过期后仍会回到账号密码体验。\n5. 下载进度不够连续，选择后立即回到主界面。", accent=C.red)
    add_card(slide, 4.82, 2.65, 3.7, 3.55, "已有基础", "1. AuthSession 已进入 Keychain，可恢复会话。\n2. 登录后会自动同步模型清单。\n3. 已有 download-ticket、加密包、SHA-256 校验。\n4. 已有设备绑定和离线租约校验。\n5. iPhone 与 platform 已通过 /api/cloud/v1 兼容。", accent=C.green)
    add_card(slide, 8.92, 2.65, 3.7, 3.55, "改造重点", "不是重写模型下载，而是重做三件事：\n1. 平台地址如何进入 App。\n2. 用户身份如何转成设备会话。\n3. 模型如何从“点下载”变成“自动就绪”。", accent=C.blue)

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "目标产品形态：平台驱动的“设备绑定 + 模型就绪”", "用户感知的是设备被组织接管，模型由平台分配；不是在手机上配置服务器和下载链接。")
    add_card(slide, 0.72, 1.62, 3.55, 1.7, "iPhone App", "默认展示相机、推理状态、模型状态。\n首次打开只出现“绑定平台”。\n高级设置里才允许手动 URL / 账号密码。", accent=C.cyan)
    add_card(slide, 4.88, 1.62, 3.55, 1.7, "vino_platform", "负责组织、用户、设备、模型授权、下载票据、审计。\n平台生成绑定邀请，不让用户输入 URL。", accent=C.blue)
    add_card(slide, 9.04, 1.62, 3.55, 1.7, "模型交付", "按 entitlement 返回可用模型。\n单模型自动下载，多模型展示选择。\n下载后自动安装、激活、续租。", accent=C.green)
    add_textbox(slide, 0.82, 4.05, 11.7, 0.35, "一句话形态", size=15, bold=True, color=C.ink)
    add_textbox(slide, 0.82, 4.55, 11.7, 0.78, "平台管理员在 Web 后台授权设备；iPhone 通过扫码或深链完成一次性绑定；之后 App 自动恢复会话、同步模型、下载并激活模型。", size=22, bold=True, color=C.dark)
    add_bullets(slide, 0.92, 5.65, 11.2, 0.75, ["普通用户日常不接触 URL、token、下载地址。", "现场工程师仍可在高级设置中切换环境或排查网络。"], size=13, color=C.muted)

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "角色与界面：三类人看到三种不同产品", "不要把所有配置都堆到 iPhone 主界面；按角色拆分入口。")
    add_card(slide, 0.72, 1.55, 3.55, 4.65, "终端操作员", "看到：\n• 相机画面\n• 当前组织 / 设备名\n• 当前模型是否就绪\n• 重新同步 / 重试\n\n不看到：\n• 平台 URL\n• access token\n• 下载 ticket\n• 复杂账号体系", accent=C.green)
    add_card(slide, 4.88, 1.55, 3.55, 4.65, "平台管理员", "看到：\n• 绑定 iPhone 按钮\n• 设备列表、在线/最后同步\n• 分配模型授权\n• 下载、续租、失败审计\n\n负责：\n• 管设备\n• 管账号\n• 管模型权限", accent=C.blue)
    add_card(slide, 9.04, 1.55, 3.55, 4.65, "现场工程师", "看到：\n• 高级设置\n• 手动平台 URL\n• 账号密码登录\n• 网络诊断\n• 日志导出\n\n用途：\n• 私有化交付\n• ECS/反代调试\n• 演示环境切换", accent=C.amber)

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "主操作流程：二维码 / 深链绑定", "这是推荐主路径，适合平台已经上线到阿里云 ECS 并有固定 HTTPS 域名的情况。")
    add_flow(slide, 0.62, 1.5, ["平台登录", "生成绑定码", "iPhone 扫码", "兑换设备会话", "同步模型", "自动下载", "模型就绪"], colors=[C.blue, C.blue, C.cyan, C.green, C.cyan, C.amber, C.green], box_w=1.32, box_h=0.58, gap=0.24)
    add_numbered(
        slide,
        0.78,
        2.52,
        5.85,
        3.8,
        [
            "管理员在平台设备页点击“绑定 iPhone”。",
            "平台生成一次性 invite code，并渲染二维码或复制链接。",
            "iPhone 打开 App，点“绑定平台”，扫码或打开 Universal Link。",
            "App 使用 code + deviceId + deviceName 兑换设备 session。",
            "平台登记设备，返回组织、用户、权限、baseURL 和 token。",
            "App 自动拉取模型清单，若只有一个模型则直接下载并激活。",
        ],
        size=12,
    )
    add_card(slide, 7.1, 2.48, 5.25, 1.2, "二维码内容原则", "只包含短期 code 和平台入口，不包含用户密码、不包含长期 token、不包含模型真实下载地址。", accent=C.red)
    add_card(slide, 7.1, 3.96, 5.25, 1.2, "绑定成功后的日常体验", "后续打开 App 自动恢复 Keychain 里的设备会话，后台同步租约和模型状态。", accent=C.green)
    add_card(slide, 7.1, 5.44, 5.25, 1.0, "必要兜底", "扫码失败时使用 6 位设备码，在平台端输入确认。", accent=C.amber)

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "iPhone 端产品流程", "从“配置型 UI”改成“状态型 UI”：让用户知道当前能不能推理，以及下一步是什么。")
    add_table_like(
        slide,
        0.7,
        1.48,
        [
            ["场景", "主界面状态", "用户动作", "系统动作"],
            ["首次打开", "未绑定平台", "点击绑定平台", "展示扫码 / 打开链接 / 输入设备码"],
            ["绑定成功", "已连接组织", "无需操作", "保存 session，自动同步模型清单"],
            ["有 1 个模型", "正在准备模型", "等待", "申请 ticket，下载，校验，安装，激活"],
            ["有多个模型", "可选择模型", "选择目标模型", "下载并激活所选模型"],
            ["离线可用", "模型已就绪", "直接推理", "按离线租约运行，后台等待续租"],
            ["会话过期", "需要重新授权", "重新扫码/设备码", "不要求输入 URL，重新兑换设备会话"],
        ],
        widths=[1.55, 2.0, 2.0, 5.8],
        row_h=0.58,
        header=True,
    )
    add_card(slide, 0.72, 6.2, 5.8, 0.75, "UI 调整", "主相机 Overlay 隐藏登录网址、账号、密码；只显示平台连接状态、模型状态、同步/重试按钮。", accent=C.blue)
    add_card(slide, 6.85, 6.2, 5.8, 0.75, "高级设置", "保留手动 URL + 账号密码，用于开发、演示、ECS 首次排查和私有化现场。", accent=C.amber)

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "平台端产品流程", "平台负责把“人、组织、设备、模型授权”串起来，iPhone 只是执行被授权的模型同步。")
    add_card(slide, 0.72, 1.55, 3.7, 1.4, "1. 设备绑定入口", "设备页增加“绑定 iPhone”。选择组织、操作者、有效期和备注，生成二维码。", accent=C.blue)
    add_card(slide, 4.82, 1.55, 3.7, 1.4, "2. 权限分配", "可以在生成 invite 前预分配模型，也可以先绑定设备后分配 entitlement。", accent=C.green)
    add_card(slide, 8.92, 1.55, 3.7, 1.4, "3. 状态监控", "展示最后上线时间、当前模型、租约到期、下载失败原因、设备是否被封禁。", accent=C.cyan)
    add_card(slide, 0.72, 3.35, 3.7, 1.4, "4. 审计与撤销", "记录 invite 创建、兑换、模型下载、租约续期；支持吊销设备会话和设备授权。", accent=C.red)
    add_card(slide, 4.82, 3.35, 3.7, 1.4, "5. 现场交付", "ECS 公网 HTTPS 地址由平台配置；二维码自动携带正确入口，避免用户输入。", accent=C.amber)
    add_card(slide, 8.92, 3.35, 3.7, 1.4, "6. 运维诊断", "设备页提供 requestId、错误码、最近 20 次同步/下载事件，便于定位生产问题。", accent=C.slate)
    add_textbox(slide, 0.82, 5.35, 11.6, 0.55, "平台端最重要的变化：把“登录 iPhone”改成“授权一台 iPhone 设备”。", size=20, bold=True, color=C.dark)

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "API 形态：在现有兼容 API 上增加绑定层", "保留现有登录、模型清单、下载票据；新增一次性 invite claim，不破坏当前 iPhone 下载实现。")
    add_table_like(
        slide,
        0.72,
        1.48,
        [
            ["阶段", "接口", "状态", "说明"],
            ["创建绑定邀请", "POST /api/platform/v1/device-invites", "新增", "平台管理员创建一次性 code，设置 TTL、组织和可选模型权限"],
            ["兑换设备会话", "POST /api/cloud/v1/device-invites/{code}/claim", "新增", "iPhone 提交 deviceId/deviceName/platform，返回 AuthSession"],
            ["恢复/验证会话", "GET /api/cloud/v1/me 或 /session", "建议新增", "App 启动时验证 Keychain session 是否仍有效"],
            ["模型清单", "GET /api/cloud/v1/models", "已有", "只返回当前设备/用户有权使用的模型"],
            ["下载票据", "POST /api/cloud/v1/models/{id}/download-ticket", "已有", "继续生成短期 ticket 和加密下载包"],
            ["下载模型", "GET /api/cloud/v1/download/{ticketId}", "已有", "继续走流式下载、加密、hash 校验"],
            ["租约续期", "POST /api/cloud/v1/licenses/lease/renew", "已有", "后台定时续期或打开 App 时续期"],
        ],
        widths=[1.55, 3.8, 1.15, 5.2],
        row_h=0.52,
        header=True,
    )
    add_card(slide, 0.72, 6.05, 5.8, 0.85, "兼容策略", "账号密码登录先保留，但降级为高级入口；新用户默认走 invite claim。", accent=C.amber)
    add_card(slide, 6.82, 6.05, 5.8, 0.85, "命名建议", "接口名可改，但必须表达一次性、可审计、可吊销、绑定设备四个语义。", accent=C.blue)

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "模型下载状态机", "用户不应看到下载 ticket，但产品要给出明确状态、进度和可恢复动作。")
    add_flow(slide, 0.65, 1.52, ["未绑定", "已绑定", "清单同步", "下载中", "安装中", "已激活", "租约续期"], colors=[C.red, C.blue, C.cyan, C.amber, C.amber, C.green, C.green], box_w=1.35, box_h=0.58, gap=0.25)
    add_card(slide, 0.72, 2.58, 3.75, 3.55, "自动下载规则", "1. 没有模型：显示“当前账号暂无模型”。\n2. 只有一个模型：绑定后自动下载。\n3. 多个模型：展示模型中心，默认推荐最新或管理员指定模型。\n4. 已安装同版本：不重复下载，只续租。\n5. 新版本可用：提示更新，可配置自动更新。", accent=C.green)
    add_card(slide, 4.82, 2.58, 3.75, 3.55, "下载中 UI", "展示模型名、版本、大小、下载进度、当前步骤。\n步骤包括：申请授权、下载、校验、解密、安装、激活。\n失败时给出重试，不让用户重新输入 URL。", accent=C.blue)
    add_card(slide, 8.92, 2.58, 3.75, 3.55, "恢复策略", "网络中断：重新申请 ticket 后续传或重下。\n安装失败：清理 staging 文件并重试。\ntoken 失效：重新扫码或设备码授权。\n设备被封禁：展示联系管理员。", accent=C.red)

    # Slide 10
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "异常与兜底流程", "把异常留给系统处理，用户只看到可执行的下一步。")
    add_table_like(
        slide,
        0.72,
        1.45,
        [
            ["异常", "用户看到", "系统处理", "平台记录"],
            ["二维码过期", "绑定码已过期，请重新生成", "拒绝 claim，不保留 session", "invite.expired"],
            ["设备数超限", "设备名额已满，请联系管理员", "不发放 session", "device_limit_reached"],
            ["账号/组织无模型", "当前账号暂无可用模型", "保持绑定状态，可继续同步", "catalog.empty"],
            ["下载 ticket 过期", "下载授权已刷新，正在重试", "重新申请 ticket", "ticket.refresh"],
            ["网络失败", "网络不可用，稍后自动重试", "指数退避，保留队列", "download.retry"],
            ["token 过期", "需要重新授权设备", "清理 session，打开绑定页", "session.expired"],
            ["设备被封禁", "设备已停用", "禁止下载和续租", "device.blocked"],
        ],
        widths=[1.7, 3.0, 3.7, 3.0],
        row_h=0.52,
        header=True,
    )
    add_card(slide, 0.72, 6.15, 11.9, 0.72, "兜底入口", "保留“设备码绑定”：iPhone 展示 6 位码，管理员在平台输入确认；适合摄像头权限、二维码扫描或深链不可用场景。", accent=C.amber)

    # Slide 11
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "安全与合规边界", "体验变“无感”，但授权边界必须更清楚。")
    add_card(slide, 0.72, 1.52, 3.75, 3.95, "绑定安全", "• invite code 短期有效，建议 5-10 分钟。\n• code 只能兑换一次。\n• code 绑定创建人、组织、目标角色。\n• claim 时记录 deviceId、设备名、IP、requestId。\n• 二维码不放密码、不放长期 token。", accent=C.blue)
    add_card(slide, 4.82, 1.52, 3.75, 3.95, "模型安全", "• 继续使用 device-bound entitlement。\n• 继续使用短期 download-ticket。\n• 下载包继续加密和 hash 校验。\n• lease 到期后离线推理不可继续。\n• 平台可封禁设备并吊销租约。", accent=C.green)
    add_card(slide, 8.92, 1.52, 3.75, 3.95, "部署要求", "• 阿里云 ECS 前建议 Nginx/SLB HTTPS。\n• VINO_EXTERNAL_BASE_URL 必须是公网 HTTPS。\n• 登录和 invite claim 需要限频。\n• 生产关闭 demo seed 账号。\n• 关键事件进入 audit log。", accent=C.red)
    add_textbox(slide, 0.85, 5.95, 11.5, 0.48, "关键判断：无感不等于匿名下载；正确形态是“平台显式授权一次，终端长期无感续期”。", size=18, bold=True, color=C.dark)

    # Slide 12
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "开发计划：小步上线，不重写已有下载链路", "先把用户路径改顺，再补平台设备绑定和自动下载。")
    add_table_like(
        slide,
        0.72,
        1.45,
        [
            ["阶段", "目标", "iPhone 改动", "platform 改动", "验收"],
            ["P0", "隐藏手动配置", "默认 baseURL 走构建配置；主界面隐藏 URL/账号/密码", "文档补 ECS HTTPS 配置", "真机首次打开不要求 URL"],
            ["P1", "设备绑定骨架", "支持 deep link / QR claim，Keychain 保存设备 session", "新增 invite create/claim，设备审计", "扫码后自动登录"],
            ["P2", "模型自动就绪", "绑定后自动同步；单模型自动下载；下载进度页", "模型清单返回推荐模型/版本", "绑定后模型可直接推理"],
            ["P3", "稳定性与兜底", "设备码、重试、续租、日志导出", "设备页状态、封禁、续租审计", "异常可恢复且可追踪"],
        ],
        widths=[0.85, 1.85, 3.4, 3.25, 2.35],
        row_h=0.72,
        header=True,
    )
    add_card(slide, 0.72, 5.7, 3.75, 1.0, "建议先做", "P0 + P1 骨架。收益最大，风险最小，不碰模型包格式。", accent=C.green)
    add_card(slide, 4.82, 5.7, 3.75, 1.0, "不建议先做", "不要先引入复杂 SSO、对象存储、多租户支付闭环；会拖慢终端体验验证。", accent=C.red)
    add_card(slide, 8.92, 5.7, 3.75, 1.0, "并行事项", "确定平台公网域名、HTTPS 证书、App deep link / Universal Link 策略。", accent=C.blue)

    # Slide 13
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "验收标准与最终用户路径", "这套改造是否成功，用可观测的体验指标判断。")
    add_card(slide, 0.72, 1.55, 5.7, 2.0, "用户路径验收", "1. 普通用户无需输入 URL。\n2. 普通用户无需输入账号密码。\n3. 扫码绑定后自动进入已连接状态。\n4. 有一个授权模型时自动下载并激活。\n5. token 过期不要求用户理解服务器地址。", accent=C.green)
    add_card(slide, 6.92, 1.55, 5.7, 2.0, "技术验收", "1. invite code 一次性和短 TTL 生效。\n2. session 存 Keychain，可恢复，可吊销。\n3. download-ticket 仍短期有效。\n4. 设备绑定、租约、审计日志完整。\n5. 所有失败返回可读错误和 requestId。", accent=C.blue)
    add_card(slide, 0.72, 4.05, 3.75, 1.55, "体验指标", "首次绑定 < 60 秒\n单模型就绪 < 2 次点击\n下载失败可一键重试", accent=C.cyan)
    add_card(slide, 4.82, 4.05, 3.75, 1.55, "运营指标", "设备绑定成功率\n模型下载成功率\n续租失败率\n设备封禁处理时长", accent=C.amber)
    add_card(slide, 8.92, 4.05, 3.75, 1.55, "上线判断", "P0/P1 后可开始小范围现场试用；P2 完成后再作为默认终端体验。", accent=C.slate)
    add_textbox(slide, 0.86, 6.25, 11.6, 0.38, "最终形态：平台授权设备，iPhone 自动拿模型，用户只关心拍摄和推理是否可用。", size=19, bold=True, color=C.dark)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
